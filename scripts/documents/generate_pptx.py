#!/usr/bin/env python3
"""Generate PPTX from a structured markdown file.

Usage:
    python scripts/documents/generate_pptx.py [input.md] [output.pptx]

Markdown format:
    # Title          → slide title (cyan, large)
    ## Subtitle      → subtitle (gray, below title)
    - Item           → bullet point (white)
    - **Bold item**  → bold bullet (cyan, acts as sub-header)
      - Sub-item     → indented bullet (gray, level 1)
    > Blockquote     → footer note (orange, bottom of slide)
    ---              → slide separator
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Theme
BG = RGBColor(0x11, 0x18, 0x27)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
GREEN = RGBColor(0x34, 0xD3, 0x99)

BOLD_RE = re.compile(r'^\*\*(.+?)\*\*:?\s*(.*)')


def parse_md(path):
    """Parse markdown into list of slide dicts."""
    slides = []
    current = None

    for raw_line in Path(path).read_text(encoding='utf-8').splitlines():
        line = raw_line.rstrip()

        if line == '---':
            if current:
                slides.append(current)
            current = None
            continue

        if current is None:
            current = {'title': '', 'subtitle': '', 'bullets': [], 'footers': []}

        if line.startswith('# ') and not line.startswith('## '):
            current['title'] = line[2:].strip()
        elif line.startswith('## '):
            current['subtitle'] = line[3:].strip()
        elif line.startswith('> '):
            current['footers'].append(line[2:].strip())
        elif line.startswith('  - '):
            # indented bullet (level 1)
            text = line[4:].strip()
            bold_m = BOLD_RE.match(text)
            if bold_m:
                label = bold_m.group(1)
                rest = bold_m.group(2)
                current['bullets'].append({'text': f"{label}: {rest}" if rest else label, 'level': 1, 'bold': True, 'color': 'cyan'})
            else:
                current['bullets'].append({'text': text, 'level': 1, 'bold': False, 'color': 'gray'})
        elif line.startswith('- '):
            text = line[2:].strip()
            bold_m = BOLD_RE.match(text)
            if bold_m:
                label = bold_m.group(1)
                rest = bold_m.group(2)
                current['bullets'].append({'text': f"{label}: {rest}" if rest else label, 'level': 0, 'bold': True, 'color': 'cyan'})
            else:
                current['bullets'].append({'text': text, 'level': 0, 'bold': False, 'color': 'white'})
        # skip empty lines and other content

    if current:
        slides.append(current)

    return slides


def render_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

    has_bullets = bool(slide_data['bullets'])
    has_footer = bool(slide_data['footers'])

    # Title
    if slide_data['title']:
        is_title_only = not has_bullets and not slide_data['subtitle']
        top = 2.5 if is_title_only else 0.4
        size = 40 if is_title_only else 32
        align = PP_ALIGN.CENTER if is_title_only else PP_ALIGN.LEFT
        left = 1 if is_title_only else 0.8

        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data['title']
        p.font.size = Pt(size)
        p.font.color.rgb = CYAN
        p.font.bold = True
        p.alignment = align

    # Subtitle
    if slide_data['subtitle']:
        is_title_only = not has_bullets
        top = 3.4 if is_title_only else 1.0
        align = PP_ALIGN.CENTER if is_title_only else PP_ALIGN.LEFT
        left = 1 if is_title_only else 0.8
        color = GRAY

        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(8.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data['subtitle']
        p.font.size = Pt(18 if is_title_only else 16)
        p.font.color.rgb = color
        p.alignment = align

    # Bullets
    if has_bullets:
        bullet_top = 1.7 if slide_data['subtitle'] else 1.2
        max_height = 5.0 if not has_footer else 4.3

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(bullet_top), Inches(8.5), Inches(max_height))
        tf = txBox.text_frame
        tf.word_wrap = True

        colors = {'cyan': CYAN, 'white': WHITE, 'gray': GRAY, 'green': GREEN}

        for i, bullet in enumerate(slide_data['bullets']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = bullet['text']
            p.font.size = Pt(16 if bullet['bold'] and bullet['level'] == 0 else 14)
            p.font.color.rgb = colors.get(bullet['color'], WHITE)
            p.font.bold = bullet['bold']
            p.level = bullet['level']
            p.space_before = Pt(6 if bullet['bold'] else 3)

    # Footers
    if has_footer:
        footer_top = 5.5
        for footer_text in slide_data['footers']:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(footer_top), Inches(8.5), Inches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = footer_text
            p.font.size = Pt(13)
            p.font.color.rgb = ORANGE
            footer_top += 0.5


def main():
    default_input = Path(__file__).parent / 'presentation_content.md'
    default_output = Path(__file__).parent.parent.parent / 'documents' / 'Mabinogi_Marketplace_Presentation.pptx'

    input_path = sys.argv[1] if len(sys.argv) > 1 else str(default_input)
    output_path = sys.argv[2] if len(sys.argv) > 2 else str(default_output)

    slides = parse_md(input_path)
    print(f"Parsed {len(slides)} slides from {input_path}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        render_slide(prs, slide_data)

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == '__main__':
    main()
